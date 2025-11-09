import openai
import PyPDF2
from pptx import Presentation
import speech_recognition as sr
import json
import re
from datetime import datetime
import os

class PresentationEvaluator:
    def __init__(self, openai_client):
        self.client = openai_client
        self.recognizer = sr.Recognizer()
        
    def extract_presentation_content(self, file_path):
        """Extract text content from PPT/PPTX or PDF files"""
        content = []
        file_extension = file_path.lower().split('.')[-1]
        
        try:
            if file_extension in ['ppt', 'pptx']:
                content = self._extract_from_powerpoint(file_path)
            elif file_extension == 'pdf':
                content = self._extract_from_pdf(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")
                
            return content
            
        except Exception as e:
            raise Exception(f"Error extracting content: {str(e)}")
    
    def _extract_from_powerpoint(self, file_path):
        """Extract text from PowerPoint slides"""
        slides_content = []
        
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_content.append({
                        'slide_number': i + 1,
                        'content': '\n'.join(slide_text)
                    })
                    
        except Exception as e:
            raise Exception(f"Error reading PowerPoint file: {str(e)}")
            
        return slides_content
    
    def _extract_from_pdf(self, file_path):
        """Extract text from PDF pages"""
        pages_content = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        pages_content.append({
                            'slide_number': i + 1,
                            'content': text
                        })
                        
        except Exception as e:
            raise Exception(f"Error reading PDF file: {str(e)}")
            
        return pages_content
    
    def analyze_presentation(self, content):
        """Analyze presentation content and generate evaluation criteria"""
        try:
            # Combine all content for analysis
            full_text = '\n'.join([slide['content'] for slide in content])
            
            # Use OpenAI to analyze the presentation
            prompt = f"""
            Analyze this presentation content and provide a structured analysis:
            
            Content:
            {full_text}
            
            Please provide:
            1. Main topic/theme
            2. Key concepts covered
            3. Technical algorithms/methodologies mentioned
            4. Project complexity level
            5. Suggested evaluation questions
            6. Areas where students might need deeper explanation
            
            Format as JSON with keys: topic, concepts, algorithms, complexity, suggested_questions, focus_areas
            """
            
            response = self.client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert presentation evaluator analyzing technical content."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3
            )
            
            analysis = json.loads(response.choices[0].message.content)
            
            # Add slide-specific analysis
            analysis['slide_count'] = len(content)
            analysis['slides'] = content
            
            return analysis
            
        except Exception as e:
            raise Exception(f"Error analyzing presentation: {str(e)}")
    
    def speech_to_text(self, audio_file_path):
        """Convert speech audio to text"""
        try:
            with sr.AudioFile(audio_file_path) as source:
                audio = self.recognizer.record(source)
                
            # Using Google Speech Recognition
            text = self.recognizer.recognize_google(audio)
            return text
            
        except sr.UnknownValueError:
            return "Could not understand audio"
        except sr.RequestError as e:
            return f"Error with speech recognition: {str(e)}"
        except Exception as e:
            return f"Error processing audio: {str(e)}"
    
    def generate_response(self, transcript, session_id):
        """Generate AI response or question based on student's speech"""
        try:
            # Get session context (would need database access)
            # For now, generate contextual responses
            
            prompt = f"""
            You are evaluating a student's presentation. The student just said:
            "{transcript}"
            
            Based on this, provide one of the following:
            1. An encouraging acknowledgment if they're doing well
            2. A clarifying question if something needs more detail
            3. A probing question to test deeper understanding
            4. Guidance if they seem stuck
            
            Be supportive but challenging. Keep responses concise and conversational.
            Focus on testing understanding of algorithms, concepts, and methodology.
            """
            
            response = self.client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert technical presentation evaluator having a conversation with a student."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=150
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            return f"I understand. Could you please elaborate more on that point?"
    
    def calculate_final_scores(self, evaluation_data, session_data):
        """Calculate final scores based on evaluation criteria"""
        try:
            # Get analysis and transcript
            analysis = session_data.get('analysis', {})
            transcript = evaluation_data.get('transcript', [])
            questions_asked = evaluation_data.get('questions_asked', [])
            answers_given = evaluation_data.get('answers_given', [])
            
            # Combine all student speech for analysis
            full_transcript = ' '.join(transcript)
            
            prompt = f"""
            Evaluate this presentation based on the following criteria:
            
            Student's Full Presentation Transcript:
            {full_transcript}
            
            Questions Asked: {questions_asked}
            Answers Given: {answers_given}
            
            Presentation Analysis: {json.dumps(analysis)}
            
            Provide scores (out of the maximum points indicated):
            1. Project Content (20 points) - Depth, relevance, correctness of the topic
            2. Algorithm Used (15 points) - Choice and justification of ML algorithms
            3. Student Skill Level (15 points) - Mastery of concepts, confidence, and critical thinking
            4. Slide Design & Visuals (10 points) - Clarity, aesthetics, and information delivery
            5. Communication & Delivery (20 points) - Oral presentation skill, engagement, flow
            6. Handling of Questions (10 points) - Ability to answer, clarity, and adaptability
            7. Research Process & Methodology (10 points) - Approach, application, reproducibility
            
            Also provide detailed feedback for each category and overall suggestions for improvement.
            
            Format as JSON with keys: scores (object with category names as keys), feedback (object with category names as keys), total_score
            """
            
            response = self.client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert evaluator providing detailed, fair, and constructive feedback."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3
            )
            
            result = json.loads(response.choices[0].message.content)
            
            # Ensure scores are within limits
            max_scores = {
                'Project Content': 20,
                'Algorithm Used': 15,
                'Student Skill Level': 15,
                'Slide Design & Visuals': 10,
                'Communication & Delivery': 20,
                'Handling of Questions': 10,
                'Research Process & Methodology': 10
            }
            
            scores = result.get('scores', {})
            for category, max_score in max_scores.items():
                if category in scores:
                    scores[category] = min(max_score, max(0, scores[category]))
                else:
                    scores[category] = 0
            
            result['scores'] = scores
            result['total_score'] = sum(scores.values())
            
            return result
            
        except Exception as e:
            # Return default scores if AI evaluation fails
            default_scores = {
                'Project Content': 10,
                'Algorithm Used': 7,
                'Student Skill Level': 7,
                'Slide Design & Visuals': 5,
                'Communication & Delivery': 10,
                'Handling of Questions': 5,
                'Research Process & Methodology': 5
            }
            
            return {
                'scores': default_scores,
                'total_score': sum(default_scores.values()),
                'feedback': {'error': f'AI evaluation failed: {str(e)}'}
            }